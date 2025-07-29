import os
import base64
import json
from typing import List, Dict, Any
from openai import OpenAI
import logging
import concurrent.futures
import PyPDF2
from dotenv import load_dotenv
import numpy as np
import cv2

load_dotenv()

logger = logging.getLogger(__name__)

# Configure logging
logging.basicConfig(
    level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s"
)

load_dotenv()

# Instead of initializing at import, create a function to get the client
client = None


def get_openai_client():
    """Get an OpenAI client with proper configuration"""
    try:
        api_key = os.environ.get("OPENAI_API_KEY")
        if not api_key:
            logger.error("OPENAI_API_KEY environment variable is not set")
            raise ValueError("OPENAI_API_KEY environment variable is not set")

        logger.debug("Attempting to create OpenAI client...")

        # Create client with explicit parameters
        client = OpenAI(api_key=api_key, timeout=60.0, max_retries=3)
        logger.debug(
            "OpenAI client created, now testing..."
        )  # Log after client creation

        # Test the client by listing available models
        models = client.models.list()
        if models and hasattr(models, "data") and len(models.data) > 0:
            logger.debug(
                f"OpenAI client initialized successfully, available models: {len(models.data)}"
            )
            return client
        else:
            logger.error("OpenAI client initialized but returned no models.")
            return None

    except Exception as e:
        logger.error(f"Error initializing OpenAI client: {str(e)}")
        return None


def encode_image_to_base64(image_path):
    """
    Encode an image to base64.

    Args:
        image_path: Path to the image file or numpy array

    Returns:
        Base64 encoded string of the image
    """
    logger.debug("Encoding image to base64")

    # Handle numpy array (from CV2)
    if isinstance(image_path, np.ndarray):
        success, buffer = cv2.imencode(".jpg", image_path)
        if not success:
            raise ValueError("Failed to encode image")
        return base64.b64encode(buffer).decode("utf-8")

    # Handle file path
    with open(image_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode("utf-8")


def create_extraction_prompt(schema: Dict[str, Any], page_count: int) -> str:
    """
    Create a prompt instructing the model how to extract data according to the schema.

    Args:
        schema: JSON schema defining the structure
        page_count: Number of pages in the document

    Returns:
        Prompt string for the model
    """
    # Create a compact JSON representation of the schema
    schema_str = json.dumps(schema, indent=2)

    prompt = f"""
    You are an expert at extracting structured data from documents.
    
    I have a document with {page_count} pages that has been converted to images. I need you to extract specific information from these images according to the following JSON schema:
    
    {schema_str}
    
    Please follow these instructions carefully:
    
    1. Examine all {page_count} images thoroughly to find the requested information.
    2. Extract the exact text from the document that matches each field in the schema.
    3. If you cannot find information for a specific field in any of the pages, return an empty string or null value for that field.
    4. For array fields, include all instances found throughout the document.
    5. Maintain the structure defined in the schema exactly.
    6. Return only the extracted data as a valid JSON object, matching the structure of the schema.
    7. Do not add any explanatory text or notes outside the JSON structure.
    8. Be precise and accurate in your extraction.
    9. If text is unclear or ambiguous, make your best guess based on context.
    10. For dates, numbers, and other formatted data, maintain the format as shown in the document.
    11. IMPORTANT: Your response MUST be a valid JSON object that exactly matches the structure of the provided schema.
    12. IMPORTANT: Do not include any explanations, just return the JSON object.
    
    Your response should be a valid JSON object containing only the extracted data.
    """

    return prompt


def semantic_chunk_with_structured_output(text: str) -> List[Dict[str, Any]]:
    """
    Use OpenAI API to create semantic chunks from text using JSON mode.

    Args:
        text: The text to chunk

    Returns:
        List of chunks with metadata
    """

    # If text is too long, split it first using a simpler method
    # and then process each part in parallel
    if len(text) > 25000:
        logger.info(
            "Text too long for direct semantic chunking, applying parallel processing"
        )
        return process_long_text_semantically(text)

    try:
        # Get the OpenAI client
        openai_client = get_openai_client()

        # Create a prompt for the OpenAI model with JSON mode
        response = openai_client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {
                    "role": "system",
                    "content": "You are an expert at analyzing and dividing text into meaningful semantic chunks. Your output should be valid JSON.",
                },
                {
                    "role": "user",
                    "content": f"""Please analyze the following text and divide it into logical semantic chunks. 
                    Each chunk should represent a cohesive unit of information or a distinct section.
                    
                    Return your results as a JSON object with this structure:
                    {{
                        "chunks": [
                            {{
                                "text": "the text of the chunk",
                                "title": "a descriptive title for this chunk",
                                "position": "beginning/middle/end"
                            }},
                            ...
                        ]
                    }}
                    
                    Text to chunk:
                    
                    {text}""",
                },
            ],
            max_tokens=4000,
            temperature=0.1,
            response_format={"type": "json_object"},
        )

        # Parse the response
        result = json.loads(response.choices[0].message.content)

        # Convert the response to our standard chunk format
        chunks = []
        current_position = 0

        for i, chunk_data in enumerate(result.get("chunks", [])):
            chunk_text = chunk_data.get("text", "")
            # Find the chunk in the original text to get accurate character positions
            start_position = text.find(chunk_text, current_position)
            if start_position == -1:
                # If exact match not found, use approximate position
                start_position = current_position

            end_position = start_position + len(chunk_text)

            chunks.append(
                {
                    "text": chunk_text,
                    "metadata": {
                        "title": chunk_data.get("title", f"Chunk {i + 1}"),
                        "position": chunk_data.get("position", "unknown"),
                        "start_char": start_position,
                        "end_char": end_position,
                        "strategy": "semantic",
                    },
                }
            )

            current_position = end_position

        return chunks

    except Exception as e:
        logger.error(f"Error in semantic chunking with JSON mode: {str(e)}")
        # Fall back to paragraph chunking if semantic chunking fails
        logger.info("Falling back to paragraph chunking")
        # We'll just do basic paragraph chunking here
        paragraphs = text.split("\n\n")
        paragraphs = [p.strip() for p in paragraphs if p.strip()]

        chunks = []
        current_position = 0

        for i, paragraph in enumerate(paragraphs):
            start_position = text.find(paragraph, current_position)
            if start_position == -1:
                start_position = current_position

            end_position = start_position + len(paragraph)

            chunks.append(
                {
                    "text": paragraph,
                    "metadata": {
                        "title": f"Paragraph {i + 1}",
                        "position": "unknown",
                        "start_char": start_position,
                        "end_char": end_position,
                        "strategy": "paragraph",  # Fall back strategy
                    },
                }
            )

            current_position = end_position

        return chunks


def process_long_text_semantically(text: str) -> List[Dict[str, Any]]:
    """
    Process a long text by breaking it into smaller pieces and chunking each piece semantically.
    Uses parallel processing and JSON mode for better performance.

    Args:
        text: The long text to process

    Returns:
        List of semantic chunks
    """
    # Create chunks of 25000 characters with 500 character overlap
    text_chunks = []
    chunk_size = 25000
    overlap = 500
    start = 0
    text_length = len(text)

    while start < text_length:
        end = min(start + chunk_size, text_length)
        text_chunks.append(text[start:end])
        start = end - overlap if end < text_length else text_length

    # Process each chunk in parallel
    all_semantic_chunks = []
    with concurrent.futures.ThreadPoolExecutor() as executor:
        # Define a worker function
        def process_chunk(chunk_text):
            try:
                # Get the OpenAI client
                openai_client = get_openai_client()

                # Process this chunk with JSON mode
                response = openai_client.chat.completions.create(
                    model="gpt-4o",
                    messages=[
                        {
                            "role": "system",
                            "content": "You are an expert at analyzing and dividing text into meaningful semantic chunks. Your output should be valid JSON.",
                        },
                        {
                            "role": "user",
                            "content": f"""Please analyze the following text and divide it into logical semantic chunks. 
                            Each chunk should represent a cohesive unit of information or a distinct section.
                            
                            Return your results as a JSON object with this structure:
                            {{
                                "chunks": [
                                    {{
                                        "text": "the text of the chunk",
                                        "title": "a descriptive title for this chunk",
                                        "position": "beginning/middle/end"
                                    }},
                                    ...
                                ]
                            }}
                            
                            Text to chunk:
                            
                            {chunk_text}""",
                        },
                    ],
                    max_tokens=4000,
                    temperature=0.1,
                    response_format={"type": "json_object"},
                )

                # Parse the response
                result = json.loads(response.choices[0].message.content)

                # Convert the response to our standard chunk format
                sub_chunks = []
                current_position = 0

                for i, chunk_data in enumerate(result.get("chunks", [])):
                    chunk_text = chunk_data.get("text", "")
                    # Find position in the original chunk
                    start_position = chunk_text.find(chunk_text, current_position)
                    if start_position == -1:
                        start_position = current_position

                    end_position = start_position + len(chunk_text)

                    sub_chunks.append(
                        {
                            "text": chunk_text,
                            "metadata": {
                                "title": chunk_data.get("title", f"Subchunk {i + 1}"),
                                "position": chunk_data.get("position", "unknown"),
                                "start_char": start_position,
                                "end_char": end_position,
                                "strategy": "semantic",
                            },
                        }
                    )

                    current_position = end_position

                return sub_chunks
            except Exception as e:
                logger.error(
                    f"Error processing semantic subchunk with JSON mode: {str(e)}"
                )
                return []

        # Submit all tasks and gather results
        futures = [executor.submit(process_chunk, chunk) for chunk in text_chunks]
        for future in concurrent.futures.as_completed(futures):
            all_semantic_chunks.extend(future.result())

    return all_semantic_chunks


def extract_text_from_pdf(pdf_path: str) -> str:
    """
    Extract text from a PDF file with optimized performance.
    Uses parallel processing for multi-page PDFs.

    Args:
        pdf_path: Path to the PDF file

    Returns:
        Extracted text from the PDF
    """

    try:
        with open(pdf_path, "rb") as file:
            reader = PyPDF2.PdfReader(file)

            # Function to extract text from a page
            def extract_page_text(page_idx):
                try:
                    page = reader.pages[page_idx]
                    text = page.extract_text() or ""
                    return text
                except Exception as e:
                    logger.warning(
                        f"Error extracting text from page {page_idx}: {str(e)}"
                    )
                    return ""

            # For small PDFs, sequential processing is faster
            if len(reader.pages) <= 5:
                all_text = ""
                for i in range(len(reader.pages)):
                    all_text += extract_page_text(i) + "\n\n"
            else:
                # Process pages in parallel for larger PDFs
                with concurrent.futures.ThreadPoolExecutor() as executor:
                    results = list(
                        executor.map(extract_page_text, range(len(reader.pages)))
                    )
                all_text = "\n\n".join(results)

        return all_text
    except Exception as e:
        logger.error(f"Error extracting text from PDF: {str(e)}")
        raise


def extract_text_from_docx(docx_path: str) -> str:
    """
    Extract text from a DOCX file.

    Args:
        docx_path: Path to the DOCX file

    Returns:
        Extracted text from the DOCX
    """
    try:
        import docx  # python-docx package

        doc = docx.Document(docx_path)
        full_text = []

        # Extract text from paragraphs
        for para in doc.paragraphs:
            full_text.append(para.text)

        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    full_text.append(cell.text)

        return "\n\n".join(full_text)
    except Exception as e:
        logger.error(f"Error extracting text from DOCX: {str(e)}")
        raise


def extract_text_from_pptx(pptx_path: str) -> str:
    """
    Extract text from a PPTX file.

    Args:
        pptx_path: Path to the PPTX file

    Returns:
        Extracted text from the PPTX
    """
    try:
        from pptx import Presentation  # python-pptx package

        presentation = Presentation(pptx_path)
        full_text = []

        # Loop through slides
        for slide_number, slide in enumerate(presentation.slides, 1):
            slide_text = []
            slide_text.append(f"Slide {slide_number}")

            # Extract text from shapes (including text boxes)
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)

                # Extract text from tables
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text:
                                row_text.append(cell.text)
                        if row_text:
                            slide_text.append(" | ".join(row_text))

            full_text.append("\n".join(slide_text))

        return "\n\n".join(full_text)
    except Exception as e:
        logger.error(f"Error extracting text from PPTX: {str(e)}")
        raise
