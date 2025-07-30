from pptx import Presentation
from .charts import process_chart
from .images import (
    replace_shape_with_image,
    should_replace_shape_with_image,
)
from .loops import (
    is_loop_end,
    is_loop_start,
    process_loops,
    LOOP_START_PATTERN,
    LOOP_END_PATTERN,
)
from .paragraphs import process_paragraph
from .pptx_utils import remove_shape
from .tables import process_table_cell


def clear_loop_directives(prs):
    """
    Clear the text of all shapes that contain loop directives.

    Args:
        prs: The Presentation object
    """
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and hasattr(shape.text_frame, "text"):
                text = shape.text_frame.text.strip()
                if LOOP_START_PATTERN.search(text) or LOOP_END_PATTERN.search(text):
                    # Clear text at paragraph level to handle formatting
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.runs:
                            for run in paragraph.runs:
                                run.text = ""
                        else:
                            paragraph.text = ""


def render_pptx(template, context: dict, output, perm_user):
    """
    Render the PPTX template (a path string or a file-like object) using the provided context and save to output.
    'output' can be a path string or a file-like object. If it's a file-like object, it will be rewound after saving.
    """
    # Support template as a file path or file-like object.
    if isinstance(template, str):
        prs = Presentation(template)
    else:
        template.seek(0)
        prs = Presentation(template)

    errors = []

    # Process loops first - identify loop sections and duplicate slides
    slides_to_process = process_loops(prs, context, perm_user, errors)

    # Process all slides including duplicated ones from loops
    for slide_info in slides_to_process:
        slide = slide_info["slide"]
        slide_number = slide_info.get("slide_number", 0)
        extra_context = slide_info.get("extra_context", {})

        # Create slide context (include `extra_context`, which is where loop variables are)
        slide_context = {
            **context,
            **extra_context,
            "slide_number": slide_number,
        }

        # Add loop variable to context if present
        if "loop_var" in slide_info and "loop_item" in slide_info:
            slide_context[slide_info["loop_var"]] = slide_info["loop_item"]

        # Process the slide's shapes
        for shape in slide.shapes:
            # Skip loop directive shapes - we'll clear them later
            if is_loop_start(shape) or is_loop_end(shape):
                continue

            # Process the shape content
            process_shape_content(
                shape, slide, slide_context, slide_number, perm_user, errors
            )

    if errors:
        print("Rendering aborted due to the following errors:")
        for err in set(errors):
            print(f" - {err}")
        print("Output file not saved.")
        return None, errors

    # Save to output (file path or file-like object)
    if isinstance(output, str):
        prs.save(output)
    else:
        prs.save(output)
        output.seek(0)

    return output, None


def process_shape_content(shape, slide, context, slide_number, perm_user, errors):
    """Process the content of a shape based on its type."""
    # 1) Check if this shape should be replaced with an image.
    if should_replace_shape_with_image(shape):
        try:
            replace_shape_with_image(
                shape,
                slide,
                context=context,
                perm_user=perm_user,
            )
        except Exception as e:
            errors.append(f"Error processing image (slide {slide_number}): {e}")
        # Skip further processing for this shape.
        return

    # 2) Check if this shape should be removed (because it's a loop directive).
    if is_loop_start(shape) or is_loop_end(shape):
        remove_shape(shape)
        return

    # 3) Process text frames (non-table).
    if hasattr(shape, "text_frame"):
        for paragraph in shape.text_frame.paragraphs:
            # Merge any placeholders that are split across multiple runs.
            try:
                process_paragraph(
                    paragraph=paragraph,
                    context=context,
                    perm_user=perm_user,
                    mode="normal",  # for text frames
                )
            except Exception as e:
                errors.append(f"Error in paragraph (slide {slide_number}): {e}")

    # 4) Process tables.
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                try:
                    process_table_cell(
                        cell=cell,
                        context=context,
                        perm_user=perm_user,
                    )
                except Exception as e:
                    errors.append(f"Error in table cell (slide {slide_number}): {e}")

    # 5) Process chart spreadsheets.
    if getattr(shape, "has_chart", False):
        try:
            process_chart(
                chart=shape.chart,
                context=context,
                perm_user=perm_user,
            )
        except Exception as e:
            errors.append(f"Error in chart (slide {slide_number}): {e}")
