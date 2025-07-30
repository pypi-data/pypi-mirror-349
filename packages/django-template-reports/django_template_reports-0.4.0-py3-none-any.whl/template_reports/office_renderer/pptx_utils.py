"""
Utility functions for PPTX manipulation.
"""

from __future__ import annotations

from typing import TYPE_CHECKING, Optional

from copy import deepcopy

if TYPE_CHECKING:
    from pptx.presentation import Presentation
    from pptx.slide import Slide
    from pptx.shapes.base import BaseShape


def duplicate_slide(
    pres: Presentation,
    slide: Slide,
    index: Optional[int] = None,
) -> Slide:
    """
    Create a duplicate of the given slide and return the new slide.

    This function duplicates a slide by:
    1. Getting the source slide
    2. Creating a new blank slide
    3. Copying all shapes from source to destination

    Args:
        pres: The Presentation object
        slide: The slide object to duplicate

    Returns:
        The newly created slide object
    """
    # (1) Create new slide (always appended to the end for now)
    new_slide = pres.slides.add_slide(slide.slide_layout)

    #  (2) Calculate where the duplicate should live
    sld_ids = pres.slides._sldIdLst
    src_idx = pres.slides.index(slide)  # position of original
    dest_idx = (
        len(sld_ids) + index
        if index is not None and index < 0  # negative slice logic
        else index if index is not None else src_idx + 1
    )
    dest_idx = max(0, min(dest_idx, len(sld_ids) - 1))  # clamp to bounds

    # (3) Move the slide-id element to that position
    new_sld_id = sld_ids[-1]  # id for slide we just appended
    sld_ids.remove(new_sld_id)
    sld_ids.insert(dest_idx, new_sld_id)
    # Delete all the shapes in the new slide
    for shape in new_slide.shapes:
        remove_shape(shape)

    # Deep-copy all shapes from the source slide to the new slide
    for shape in slide.shapes:
        new_el = deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

    return new_slide


def remove_shape(shape: BaseShape):
    """
    Remove a shape from a slide.

    This function removes a shape by removing its XML element from the slide's shape tree.

    Args:
        slide: The slide object containing the shape
        shape: The shape object to remove
    """
    parent = shape.element.getparent()
    if parent is not None:
        parent.remove(shape.element)
