import re
from pptx import Presentation

from .charts import get_raw_chart_data
from .paragraphs import merge_split_placeholders

# Pattern to match placeholders, e.g. "{{ some.placeholder }}"
PLACEHOLDER_PATTERN = re.compile(r"\{\{\s*(.*?)\s*\}\}")


def extract_top_level_context_keys_from_text(text: str) -> dict[str, list[str]]:
    """
    Given a text string, find all placeholders and extract the top-level context key.
    Returns a dict with:
        - simple_fields: keys without square brackets or periods in the placeholder
        - object_fields: keys where the placeholder includes a square bracket or period
    """
    simple_fields = set()
    object_fields = set()
    placeholders = PLACEHOLDER_PATTERN.findall(text)
    for ph in placeholders:
        ph = ph.strip()
        if ph:
            m = re.match(r"([^\.\[\]\|]+)", ph)
            if m:
                key = m.group(1).strip()
                if key == "now":
                    continue
                if ("." in ph) or ("[" in ph):
                    object_fields.add(key)
                else:
                    simple_fields.add(key)
    return {
        "simple_fields": sorted(simple_fields),
        "object_fields": sorted(object_fields),
    }


def extract_context_keys(template) -> dict[str, list[str]]:
    """
    Iterate through all slides, shapes, paragraphs and table cells in the PPTX (from a file path or file-like object),
    merging split placeholders, and return a dict with:
        - simple_fields: sorted list of unique simple keys
        - object_fields: sorted list of unique object keys
    """
    if isinstance(template, str):
        prs = Presentation(template)
    else:
        template.seek(0)
        prs = Presentation(template)

    simple_fields = set()
    object_fields = set()

    # Build a list of all texts on all slides.
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:

            # Process text frames.
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    merge_split_placeholders(paragraph)
                    texts.append(paragraph.text)
            # Process table cells.
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text_frame:
                            for paragraph in cell.text_frame.paragraphs:
                                merge_split_placeholders(paragraph)
                                texts.append(paragraph.text)
            # Process chart spreadsheets.
            if getattr(shape, "has_chart", False):
                raw_data = get_raw_chart_data(shape.chart)
                for col in raw_data:
                    for item in col:
                        texts.append(str(item))

    # Process all texts to extract context keys.
    for text in texts:
        keys = extract_top_level_context_keys_from_text(text)
        simple_fields.update(keys["simple_fields"])
        object_fields.update(keys["object_fields"])

    return {
        "simple_fields": sorted(simple_fields),
        "object_fields": sorted(object_fields),
    }
