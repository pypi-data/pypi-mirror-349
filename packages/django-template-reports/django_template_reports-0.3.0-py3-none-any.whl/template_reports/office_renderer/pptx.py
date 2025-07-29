from pptx import Presentation

from .charts import process_chart
from .images import (
    replace_shape_with_image,
    should_replace_shape_with_image,
)
from .paragraphs import process_paragraph
from .tables import process_table_cell


def render_pptx(template, context: dict, output, perm_user):
    """
    Render the PPTX template (a path string or a file-like object) using the provided context and save to output.
    'output' can be a path string or a file-like object. If it's a file-like object, it will be rewound after saving.
    """
    # Support template as a file path or file-like object.
    if isinstance(template, str):
        prs = Presentation(template)
    else:
        template.seek(0)
        prs = Presentation(template)

    errors = []

    for i, slide in enumerate(prs.slides):
        # Add extra context per slide.
        slide_number = i + 1
        slide_context = {
            **context,
            "slide_number": slide_number,
        }

        for shape in slide.shapes:
            # Check if this shape should be replaced with an image.
            if should_replace_shape_with_image(shape):
                try:
                    replace_shape_with_image(
                        shape,
                        slide,
                        context=slide_context,
                        perm_user=perm_user,
                    )
                except Exception as e:
                    errors.append(
                        f"Error processing image (slide {slide_number}): {e}"
                    )
                # Skip further processing for this shape
                continue

            # 1) Process text frames (non-table).
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    # Merge any placeholders that are split across multiple runs.
                    try:
                        process_paragraph(
                            paragraph=paragraph,
                            context=slide_context,
                            perm_user=perm_user,
                            mode="normal",  # for text frames
                        )
                    except Exception as e:
                        errors.append(f"Error in paragraph (slide {slide_number}): {e}")
            # 2) Process tables.
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        try:
                            process_table_cell(
                                cell=cell,
                                context=slide_context,
                                perm_user=perm_user,
                            )
                        except Exception as e:
                            errors.append(
                                f"Error in table cell (slide {slide_number}): {e}"
                            )
            # 3) Process chart spreadsheets.
            if getattr(shape, "has_chart", False):
                try:
                    process_chart(
                        chart=shape.chart,
                        context=slide_context,
                        perm_user=perm_user,
                    )
                except Exception as e:
                    errors.append(f"Error in chart (slide {slide_number}): {e}")


    if errors:
        print("Rendering aborted due to the following errors:")
        for err in set(errors):
            print(f" - {err}")
        print("Output file not saved.")
        return None, errors

    # Save to output (file path or file-like object)
    if isinstance(output, str):
        prs.save(output)
    else:
        prs.save(output)
        output.seek(0)

    return output, None
