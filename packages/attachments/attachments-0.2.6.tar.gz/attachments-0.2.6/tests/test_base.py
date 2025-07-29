import unittest
import os
import subprocess # Keep for now if any part of setup still needs it, though aiming to remove
from PIL import Image # For creating sample images
import wave # For creating dummy WAV files
import struct # For creating dummy WAV files
import pytest # Import pytest
# from typing import Type, Optional # Removing for now to simplify

# For pptx
try:
    from pptx import Presentation
    from pptx.util import Inches
    print("Successfully imported Presentation and Inches from pptx.")
except ImportError:
    Presentation, Inches = None, None # type: ignore
    print("Warning: python-pptx not installed. PPTX creation in test setup might fail.")


# For fitz (PyMuPDF)
try:
    import fitz as _fitz_module
    print("Successfully imported fitz as _fitz_module")
except ImportError:
    _fitz_module = None # type: ignore
    print("Warning: PyMuPDF (fitz) not imported. PDF creation in test setup might fail.")


# Define the path to the test data directory
# __file__ in test_base.py will point to tests/test_base.py
# So, os.path.dirname(__file__) is tests/
# And os.path.join(os.path.dirname(__file__), 'test_data') is tests/test_data
TEST_DATA_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'test_data')

# Define paths to common sample files
SAMPLE_PDF = os.path.join(TEST_DATA_DIR, 'sample.pdf')
SAMPLE_PPTX = os.path.join(TEST_DATA_DIR, 'sample.pptx')
GENERATED_MULTI_PAGE_PDF = os.path.join(TEST_DATA_DIR, 'multi_page.pdf')
SAMPLE_HTML = os.path.join(TEST_DATA_DIR, 'sample.html')
NON_EXISTENT_FILE = os.path.join(TEST_DATA_DIR, 'not_here.txt') # Should remain non-existent
SAMPLE_PNG = os.path.join(TEST_DATA_DIR, 'sample.png')
SAMPLE_JPG = os.path.join(TEST_DATA_DIR, 'sample.jpg')
SAMPLE_HEIC = os.path.join(TEST_DATA_DIR, 'sample.heic') # HEIC remains optional, checked for existence

# Audio files
AUDIO_DIR = os.path.join(TEST_DATA_DIR, 'audio') # Define audio subdir
SAMPLE_OGG = os.path.join(AUDIO_DIR, 'sample.ogg') 
SAMPLE_MP3 = os.path.join(AUDIO_DIR, 'sample.mp3')
SAMPLE_WAV = os.path.join(AUDIO_DIR, 'sample.wav') # This one is created as a valid silent WAV
SAMPLE_FLAC = os.path.join(AUDIO_DIR, 'sample.flac')
SAMPLE_M4A = os.path.join(AUDIO_DIR, 'sample.m4a')
SAMPLE_MP4_AUDIO = os.path.join(AUDIO_DIR, 'sample_audio.mp4')
SAMPLE_WEBM_AUDIO = os.path.join(AUDIO_DIR, 'sample_audio.webm')
USER_PROVIDED_WAV = os.path.join(TEST_DATA_DIR, 'sample_audio.wav') # User's main WAV for testing conversions

SAMPLE_DOCX = os.path.join(TEST_DATA_DIR, 'sample.docx')
SAMPLE_ODT = os.path.join(TEST_DATA_DIR, 'sample.odt')


# Helper to create a multi-page PDF for testing PDF indexing
def create_multi_page_pdf(path, num_pages=5):
    if os.path.exists(path) and os.path.getsize(path) > 0: # Check if not empty
        # Potentially add a more robust check if content can vary and needs to be specific
        # For now, if it exists and is not empty, assume it's okay.
        # print(f"{path} already exists and is not empty.")
        return
    if _fitz_module is None:
        print(f"PyMuPDF (_fitz_module) not available, cannot create {path}. Some PDF tests may fail or be skipped.")
        return
    try:
        doc = _fitz_module.open() # New PDF
        for i in range(num_pages):
            page = doc.new_page()
            page.insert_text((50, 72), f"This is page {i+1} of {num_pages}.")
        doc.save(path)
        doc.close()
        print(f"Created {path} with {num_pages} pages for testing.")
    except Exception as e:
        print(f"Could not create multi-page PDF {path}: {e}")

def _create_dummy_audio_file_for_fixture(file_path, file_type): # Renamed to avoid conflict if old class is temp. present
    """Creates a dummy audio file if it doesn't exist."""
    if os.path.exists(file_path) and os.path.getsize(file_path) > 0:
        # print(f"Dummy audio {file_path} already exists and is not empty.")
        return
    
    audio_dir = os.path.dirname(file_path)
    if not os.path.exists(audio_dir):
        os.makedirs(audio_dir, exist_ok=True)

    if file_type == 'wav':
        try:
            # Create a minimal, silent WAV file
            sample_rate = 44100; duration_ms = 10
            n_frames = int(sample_rate * (duration_ms / 1000.0))
            n_channels = 1; sampwidth = 2 # 16-bit
            comptype = "NONE"; compname = "not compressed"

            with wave.open(file_path, 'wb') as wf:
                wf.setnchannels(n_channels); wf.setsampwidth(sampwidth)
                wf.setframerate(sample_rate); wf.setnframes(n_frames)
                wf.setcomptype(comptype, compname)
                for _ in range(n_frames):
                    wf.writeframesraw(struct.pack('<h', 0))
            print(f"Fixture created dummy WAV file: {file_path}")
        except Exception as e:
            print(f"Fixture failed to create dummy WAV file {file_path}: {e}")
    else:
        # For other formats, create a small placeholder file
        try:
            with open(file_path, 'wb') as f:
                f.write(f"dummy content for {file_type}".encode('utf-8'))
            print(f"Fixture created dummy placeholder audio file: {file_path} (type: {file_type})")
        except Exception as e:
            print(f"Fixture failed to create dummy placeholder audio file {file_path}: {e}")


@pytest.fixture(scope="class", autouse=True)
def base_test_setup(request):
    cls = request.cls # The test class
    print(f"Running pytest fixture 'base_test_setup' for class {cls.__name__}")

    # Ensure TEST_DATA_DIR and AUDIO_DIR exist (though they should if files are being made in them)
    if not os.path.exists(TEST_DATA_DIR):
        os.makedirs(TEST_DATA_DIR, exist_ok=True)
    if not os.path.exists(AUDIO_DIR):
        os.makedirs(AUDIO_DIR, exist_ok=True)
        
    cls.test_output_dir = os.path.join(TEST_DATA_DIR, "test_outputs_fixture_base") 
    if not os.path.exists(cls.test_output_dir):
        os.makedirs(cls.test_output_dir, exist_ok=True)
            
    # Ensure sample PDF exists
    if not (os.path.exists(SAMPLE_PDF) and os.path.getsize(SAMPLE_PDF) > 0):
        if _fitz_module:
            try:
                doc = _fitz_module.open()
                page = doc.new_page()
                page.insert_text((50, 72), "Hello PDF!")
                doc.save(SAMPLE_PDF)
                doc.close()
                print(f"Fixture created {SAMPLE_PDF} for testing.")
            except Exception as e:
                 print(f"Fixture warning: Could not create {SAMPLE_PDF}: {e}")
        else:
            print(f"Fixture warning: PyMuPDF (_fitz_module) not available. Cannot create {SAMPLE_PDF}")
    
    create_multi_page_pdf(GENERATED_MULTI_PAGE_PDF, 5)
    cls.generated_multi_page_pdf_exists = os.path.exists(GENERATED_MULTI_PAGE_PDF) and os.path.getsize(GENERATED_MULTI_PAGE_PDF) > 0

    # Programmatically create sample.pptx
    if not (os.path.exists(SAMPLE_PPTX) and os.path.getsize(SAMPLE_PPTX) > 100): # Check size for basic validity
        if Presentation is None or Inches is None:
            print("Fixture: python-pptx (Presentation or Inches) not available. Cannot create sample.pptx.")
            cls.sample_pptx_exists = False
        else:
            try:
                if os.path.exists(SAMPLE_PPTX): # If it exists but is small/invalid
                    os.remove(SAMPLE_PPTX)

                prs = Presentation()
                slide_layout_title = prs.slide_layouts[0] 
                slide1 = prs.slides.add_slide(slide_layout_title)
                title1 = slide1.shapes.title
                title1.text = "Slide 1 Title"

                slide_layout_content = prs.slide_layouts[1]
                slide2 = prs.slides.add_slide(slide_layout_content)
                title2 = slide2.shapes.title
                title2.text = "Slide 2 Title"
                body2 = slide2.placeholders[1].text_frame
                body2.text = "Content for page 2"

                slide3 = prs.slides.add_slide(slide_layout_content)
                title3 = slide3.shapes.title
                title3.text = "Slide 3 Title"
                body3 = slide3.placeholders[1].text_frame
                body3.text = "Content for page 3"

                prs.save(SAMPLE_PPTX)
                print(f"Fixture created {SAMPLE_PPTX} programmatically.")
                cls.sample_pptx_exists = True
            except Exception as e:
                print(f"Fixture: Could not create or verify sample.pptx programmatically: {e}.")
                cls.sample_pptx_exists = False
    else:
        # print(f"{SAMPLE_PPTX} already exists and seems valid enough.")
        cls.sample_pptx_exists = True # Assume valid if exists and has some size

    # Verify readability if we think it exists or was just created
    if cls.sample_pptx_exists and Presentation is not None:
        try:
            Presentation(SAMPLE_PPTX)
        except Exception as e:
            print(f"Fixture warning: {SAMPLE_PPTX} could not be reliably opened by python-pptx: {e}.")
            cls.sample_pptx_exists = False # Mark as problematic

    if not (os.path.exists(SAMPLE_HTML) and os.path.getsize(SAMPLE_HTML) > 0):
        try:
            with open(SAMPLE_HTML, "w") as f:
                f.write("<html><head><title>Sample HTML</title></head><body><h1>Main Heading</h1><p>This is a paragraph with <strong>strong emphasis</strong> and <em>italic text</em>. <a href=\"http://example.com\">Example Link</a></p><ul><li>First item</li><li>Second item</li></ul><script>console.log('test');</script></body></html>")
            print(f"Fixture created a fallback {SAMPLE_HTML} as it was missing or empty.")
        except Exception as e_html_create:
            print(f"Fixture: Could not create fallback {SAMPLE_HTML}: {e_html_create}")
    cls.sample_html_exists = os.path.exists(SAMPLE_HTML) and os.path.getsize(SAMPLE_HTML) > 0


    # PNG and JPG creation
    if not (os.path.exists(SAMPLE_PNG) and os.path.getsize(SAMPLE_PNG) > 0):
        try:
            img_png = Image.new('RGB', (10, 10), color = 'red') # Make slightly larger than 1x1
            img_png.save(SAMPLE_PNG, 'PNG')
            print(f"Fixture created {SAMPLE_PNG} programmatically.")
        except Exception as e_png:
            print(f"Fixture: Could not create {SAMPLE_PNG}: {e_png}")
    cls.sample_png_exists = os.path.exists(SAMPLE_PNG) and os.path.getsize(SAMPLE_PNG) > 0
    
    if not (os.path.exists(SAMPLE_JPG) and os.path.getsize(SAMPLE_JPG) > 0):
        try:
            img_jpg = Image.new('RGB', (10, 10), color = 'blue') # Make slightly larger
            img_jpg.save(SAMPLE_JPG, 'JPEG')
            print(f"Fixture created {SAMPLE_JPG} programmatically.")
        except Exception as e_jpg:
            print(f"Fixture: Could not create {SAMPLE_JPG}: {e_jpg}")
    cls.sample_jpg_exists = os.path.exists(SAMPLE_JPG) and os.path.getsize(SAMPLE_JPG) > 0
    
    cls.sample_heic_exists = os.path.exists(SAMPLE_HEIC) and os.path.getsize(SAMPLE_HEIC) > 0


    # Warnings for missing image files (after attempting creation)
    if not cls.sample_png_exists: print(f"Fixture CRITICAL WARNING: {SAMPLE_PNG} is still missing or empty.")
    if not cls.sample_jpg_exists: print(f"Fixture CRITICAL WARNING: {SAMPLE_JPG} is still missing or empty.")
    if not cls.sample_heic_exists: print(f"Fixture WARNING: {SAMPLE_HEIC} is missing or empty. HEIC tests might skip or use fallbacks.")

    # Create dummy audio files
    _create_dummy_audio_file_for_fixture(SAMPLE_OGG, 'ogg')
    cls.sample_ogg_exists = os.path.exists(SAMPLE_OGG) and os.path.getsize(SAMPLE_OGG) > 0
    _create_dummy_audio_file_for_fixture(SAMPLE_MP3, 'mp3')
    cls.sample_mp3_exists = os.path.exists(SAMPLE_MP3) and os.path.getsize(SAMPLE_MP3) > 0
    _create_dummy_audio_file_for_fixture(SAMPLE_WAV, 'wav') # This creates a valid silent WAV
    cls.sample_wav_exists = os.path.exists(SAMPLE_WAV) and os.path.getsize(SAMPLE_WAV) > 0
    _create_dummy_audio_file_for_fixture(SAMPLE_FLAC, 'flac')
    cls.sample_flac_exists = os.path.exists(SAMPLE_FLAC) and os.path.getsize(SAMPLE_FLAC) > 0
    _create_dummy_audio_file_for_fixture(SAMPLE_M4A, 'm4a')
    cls.sample_m4a_exists = os.path.exists(SAMPLE_M4A) and os.path.getsize(SAMPLE_M4A) > 0
    _create_dummy_audio_file_for_fixture(SAMPLE_MP4_AUDIO, 'mp4')
    cls.sample_mp4_audio_exists = os.path.exists(SAMPLE_MP4_AUDIO) and os.path.getsize(SAMPLE_MP4_AUDIO) > 0
    _create_dummy_audio_file_for_fixture(SAMPLE_WEBM_AUDIO, 'webm')
    cls.sample_webm_audio_exists = os.path.exists(SAMPLE_WEBM_AUDIO) and os.path.getsize(SAMPLE_WEBM_AUDIO) > 0

    cls.user_provided_wav_exists = os.path.exists(USER_PROVIDED_WAV) and os.path.getsize(USER_PROVIDED_WAV) > 0
    if not cls.user_provided_wav_exists:
        print(f"Fixture CRITICAL WARNING: User provided WAV {USER_PROVIDED_WAV} is missing or empty. Audio conversion tests will fail/skip.")

    # For DOCX and ODT, we assume they are provided by the user, as in the original setup.
    # The setup was using files copied from Downloads. We'll just check for their existence here.
    # If they are missing, tests that specifically need them should skip.
    cls.sample_docx_exists = os.path.exists(SAMPLE_DOCX) and os.path.getsize(SAMPLE_DOCX) > 0
    if not cls.sample_docx_exists:
        print(f"Fixture WARNING: {SAMPLE_DOCX} is missing or empty. DOCX tests might fail/skip.")
    
    cls.sample_odt_exists = os.path.exists(SAMPLE_ODT) and os.path.getsize(SAMPLE_ODT) > 0
    if not cls.sample_odt_exists:
        print(f"Fixture WARNING: {SAMPLE_ODT} is missing or empty. ODT tests might fail/skip.")

    # If unittest.TestCase is still used by test classes, they might expect tearDownClass.
    # Pytest fixtures handle teardown differently (e.g. with a yield or addfinalizer).
    # For now, this fixture replicates setUpClass only.
    # If tearDownClass logic from BaseAttachmentsTest existed, it would go here or be registered.
    # def finalizer():
    #    print(f"Tearing down pytest fixture 'base_test_setup' for class {cls.__name__}")
    # request.addfinalizer(finalizer)


# Removed the old BaseAttachmentsTest class
# class BaseAttachmentsTest(unittest.TestCase):
#    ...


if __name__ == '__main__':
    # This file is not meant to be run directly as a test suite
    # but you can run its setUpClass for debugging asset creation.
    # BaseAttachmentsTest.setUpClass()
    # print(f"SAMPLE_PDF exists: {BaseAttachmentsTest.sample_pdf_exists}")
    # print(f"SAMPLE_PPTX exists: {BaseAttachmentsTest.sample_pptx_exists}")
    # print(f"SAMPLE_PNG exists: {BaseAttachmentsTest.sample_png_exists}")
    # print(f"SAMPLE_JPG exists: {BaseAttachmentsTest.sample_jpg_exists}")
    # print(f"SAMPLE_WAV exists: {BaseAttachmentsTest.sample_wav_exists}")
    # print(f"SAMPLE_DOCX exists: {BaseAttachmentsTest.sample_docx_exists}")
    pass 